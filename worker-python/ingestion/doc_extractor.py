"""Document text extractor for Executive Digest pipeline.

Supports: PDF, DOCX, PPTX, XLSX
All processing is local — no data sent to external services.

Returns a list of page/slide/sheet chunks, each as:
  { "page": int, "text": str, "source": filename }
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional

SUPPORTED_DOC_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}


def extract_text(file_path: Path) -> List[Dict[str, Any]]:
    """Extract text chunks from a document file.

    Returns list of dicts with keys: page (1-indexed), text, source.
    Raises ValueError for unsupported file types.
    """
    suffix = file_path.suffix.lower()
    if suffix not in SUPPORTED_DOC_EXTENSIONS:
        raise ValueError(f"Unsupported document type: {suffix}")

    if suffix == ".pdf":
        return _extract_pdf(file_path)
    if suffix == ".docx":
        return _extract_docx(file_path)
    if suffix == ".pptx":
        return _extract_pptx(file_path)
    if suffix == ".xlsx":
        return _extract_xlsx(file_path)
    return []


def _extract_pdf(path: Path) -> List[Dict[str, Any]]:
    try:
        import pdfplumber
        chunks = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    chunks.append({"page": i, "text": text.strip(), "source": path.name})
        return chunks
    except ImportError:
        pass

    # fallback: pypdf
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        chunks = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append({"page": i, "text": text.strip(), "source": path.name})
        return chunks
    except ImportError:
        raise ImportError("Install pdfplumber or pypdf: pip install pdfplumber")


def _extract_docx(path: Path) -> List[Dict[str, Any]]:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("Install python-docx: pip install python-docx")

    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # Group into ~30-paragraph chunks to simulate "pages"
    chunk_size = 30
    chunks = []
    for i in range(0, len(paragraphs), chunk_size):
        text = "\n".join(paragraphs[i : i + chunk_size])
        chunks.append({"page": i // chunk_size + 1, "text": text, "source": path.name})
    return chunks


def _extract_pptx(path: Path) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Install python-pptx: pip install python-pptx")

    prs = Presentation(str(path))
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            chunks.append({"page": i, "text": "\n".join(texts), "source": path.name})
    return chunks


def _extract_xlsx(path: Path) -> List[Dict[str, Any]]:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("Install openpyxl: pip install openpyxl")

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    chunks = []
    page = 1
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            chunks.append({
                "page": page,
                "text": f"[Sheet: {sheet_name}]\n" + "\n".join(rows),
                "source": path.name,
            })
            page += 1
    wb.close()
    return chunks


def validate_document(path: Path) -> bool:
    """Return True if the file is a supported document type and exists."""
    return path.exists() and path.is_file() and path.suffix.lower() in SUPPORTED_DOC_EXTENSIONS


def combine_chunks_to_text(chunks: List[Dict[str, Any]]) -> str:
    """Flatten all chunks into a single string for LLM input."""
    return "\n\n".join(
        f"[Page {c['page']}]\n{c['text']}" for c in chunks
    )
